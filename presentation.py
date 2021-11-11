import keypress_detect
import speech_detect
import pptx
import os

pressed = True

def open_presentation(s):
  current_dir = str(os.path.dirname(os.path.abspath(__file__))) + '/' + s
  try:
    pres = pptx.Presentation(s)
    return pres
  except:
    return -1


def display_text(p):
  triggered = False
  i = 0
  for slide in p.slides:
    for shape in slide.shapes:
        for paragraph in shape.text_frame.paragraphs:
            for k, run in enumerate(paragraph.runs):
                print(run.text)
                if ((k+1) == len(paragraph.runs)):
                    print("End of page.")
                    sentences = run.text.split(". ")
                    last_sentence = sentences[-1]
                    curr_x = keypress_detect.x
                    print(last_sentence)
                    # If this is a title, skip and display sentences
                    if len(sentences) != 0:
                      while curr_x == keypress_detect.x:
                        if speech_detect.determine_phrase(last_sentence, speech_detect.listen()):
                          keypress_detect.execute_keypress()